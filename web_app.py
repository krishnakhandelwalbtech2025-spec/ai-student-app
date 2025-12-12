import streamlit as st
import google.generativeai as genai
from PIL import Image
import pypdf
from youtube_transcript_api import YouTubeTranscriptApi
from pptx import Presentation
import time
from gtts import gTTS
import io

# --- 1. CONFIGURATION ---
# Securely access key from Cloud Secrets
api_key = st.secrets["AIzaSyA_6qI9720Ze19lSfSh0-ZQRYLRaFRJc-U"]
genai.configure(api_key=api_key)

genai.configure(api_key=api_key)

genai.configure(api_key=api_key)
model = genai.GenerativeModel("gemini-2.5-flash")

st.set_page_config(page_title="Ultimate Student AI", page_icon="🎓", layout="wide")

# --- 2. GAMIFICATION ENGINE (XP SYSTEM) ---
if "xp" not in st.session_state:
    st.session_state.xp = 0
if "level" not in st.session_state:
    st.session_state.level = 1

def add_xp(amount):
    st.session_state.xp += amount
    # Level up logic: Level 1 requires 100 XP, Level 2 requires 200 XP, etc.
    xp_needed = st.session_state.level * 100
    if st.session_state.xp >= xp_needed:
        st.session_state.level += 1
        st.session_state.xp -= xp_needed
        st.balloons()
        st.toast(f"🎉 LEVEL UP! You are now Level {st.session_state.level}!")
    else:
        st.toast(f"⭐ +{amount} XP")

# --- 3. SIDEBAR (PROFILE & NAV) ---
st.sidebar.title("👤 Student Profile")
st.sidebar.write(f"**Level {st.session_state.level}**")
st.sidebar.progress(st.session_state.xp / (st.session_state.level * 100))
st.sidebar.caption(f"{st.session_state.xp} / {st.session_state.level * 100} XP to next level")

st.sidebar.markdown("---")
st.sidebar.subheader("🚀 Tool Center")

# --- FOCUS TIMER ---
timer_minutes = st.sidebar.number_input("⏱️ Focus Timer (Mins):", 1, 120, 25)
if st.sidebar.button("Start Focus Session"):
    progress_bar = st.sidebar.progress(0)
    status_text = st.sidebar.empty()
    total_seconds = timer_minutes * 60
    for i in range(total_seconds):
        mins, secs = divmod(total_seconds - i, 60)
        status_text.write(f"⏳ {mins:02d}:{secs:02d}")
        progress_bar.progress((i + 1) / total_seconds)
        time.sleep(1)
    st.sidebar.success("⏰ Session Complete!")
    add_xp(50)  # Reward for focusing

st.sidebar.markdown("---")
mode = st.sidebar.radio(
    "Select Tool:",
    [
        "📊 Dashboard",
        "🧠 Mind Mapper",       # NEW
        "📚 Smart Notes",       # UPDATED (Audio)
        "📝 Adaptive Quiz", 
        "🧮 Formula Cards",
        "📄 Document Reader", 
        "💬 Chat", 
        "📺 YouTube Summarizer", 
        "🖼️ Vision", 
        "🌐 Web Search"
    ]
)

st.title(f"{mode}")

# --- 4. HELPER FUNCTIONS ---
def get_gemini_response(prompt, tools=None):
    try:
        model = genai.GenerativeModel("gemini-2.5-flash", tools=tools)
        return model.generate_content(prompt).text
    except Exception as e:
        return f"Error: {e}"

def text_to_speech(text):
    try:
        tts = gTTS(text=text, lang='en')
        audio_fp = io.BytesIO()
        tts.write_to_fp(audio_fp)
        return audio_fp
    except:
        return None

def extract_text(uploaded_file):
    text = ""
    if uploaded_file.name.endswith(".pdf"):
        reader = pypdf.PdfReader(uploaded_file)
        text = "".join([page.extract_text() for page in reader.pages])
    elif uploaded_file.name.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): text += shape.text + "\n"
    return text

# --- 5. MODES ---

# === MIND MAPPER (NEW) ===
if mode == "🧠 Mind Mapper":
    st.write("Visual Learning: Generate a Mind Map for any topic.")
    topic = st.text_input("Enter Topic:")
    if st.button("Generate Mind Map") and topic:
        with st.spinner("Brainstorming..."):
            # Request Graphviz DOT code from Gemini
            prompt = f"""
            Create a Graphviz DOT code for a Mind Map about: {topic}.
            The graph should be hierarchical and colorful.
            Only output the code inside ```dot ... ``` blocks.
            """
            response = get_gemini_response(prompt)
            
            # Clean up response to get just the DOT code
            try:
                dot_code = response.split("```dot")[1].split("```")[0]
                st.graphviz_chart(dot_code)
                add_xp(20)
            except:
                st.error("AI couldn't generate a valid diagram structure. Try a simpler topic.")

# === SMART NOTES (WITH PODCAST) ===
elif mode == "📚 Smart Notes":
    st.write("Generate structured notes + Audio Podcast.")
    input_type = st.radio("Source:", ["Type Topic", "Upload File"])
    content = ""
    if input_type == "Type Topic":
        topic = st.text_input("Enter Topic:")
        if topic: content = f"Topic: {topic}"
    else:
        f = st.file_uploader("Upload PDF/PPTX", type=["pdf", "pptx"])
        if f: content = extract_text(f)

    if st.button("Generate Notes") and content:
        with st.spinner("Writing notes..."):
            notes = get_gemini_response(f"Create student notes for:\n{content[:10000]}")
            st.markdown(notes)
            add_xp(30)
            
            # AUDIO GENERATION
            with st.spinner("Generating Podcast Audio..."):
                audio_file = text_to_speech(notes[:500]) # Limit audio length for speed
                if audio_file:
                    st.audio(audio_file, format='audio/mp3')
                    st.success("🎧 Podcast Ready!")

# === STUDENT DASHBOARD ===
elif mode == "📊 Dashboard":
    st.subheader("👋 Student Dashboard")
    if "tasks" not in st.session_state: st.session_state.tasks = []
    
    col1, col2 = st.columns(2)
    with col1:
        new_task = st.text_input("New Task:")
        if st.button("Add") and new_task:
            st.session_state.tasks.append({"task": new_task, "done": False})
            st.rerun()
            
        for i, task in enumerate(st.session_state.tasks):
            done = st.checkbox(task["task"], key=f"t{i}", value=task["done"])
            if done and not task["done"]:
                st.session_state.tasks[i]["done"] = True
                add_xp(10) # XP for tasks
                st.rerun()
    with col2:
        st.metric("Total XP", st.session_state.xp)
        st.metric("Current Level", st.session_state.level)

# === ADAPTIVE QUIZ ===
elif mode == "📝 Adaptive Quiz":
    # (Same robust quiz logic, simplified for brevity but fully functional)
    if "quiz_state" not in st.session_state:
        st.session_state.quiz_state = "setup"
        st.session_state.score = 0
        st.session_state.quiz_data = []
        st.session_state.current_q = 0

    if st.session_state.quiz_state == "setup":
        topic = st.text_input("Quiz Topic:")
        if st.button("Start") and topic:
            with st.spinner("Preparing..."):
                raw = get_gemini_response(f"Generate 5 multiple choice questions on {topic}. Format:\nQ: Question\nA) opt\nB) opt\nCorrect: Letter")
                # (Simple parser)
                try:
                    qs = raw.split("Q:")[1:]
                    st.session_state.quiz_data = [{"q": q.split("\n")[0], "opts": [l for l in q.split("\n") if l.startswith(("A)", "B)"))], "ans": [l for l in q.split("\n") if "Correct:" in l][0].split(":")[1].strip()} for q in qs]
                    st.session_state.quiz_state = "playing"
                    st.rerun()
                except: st.error("Try again.")

    elif st.session_state.quiz_state == "playing":
        q = st.session_state.quiz_data[st.session_state.current_q]
        st.write(q["q"])
        choice = st.radio("Options", q["opts"], key=f"q{st.session_state.current_q}")
        if st.button("Submit"):
            if choice.split(")")[0] == q["ans"]:
                st.success("Correct!")
                st.session_state.score += 1
                add_xp(10) # XP for correct answer
            else: st.error(f"Wrong. Answer: {q['ans']}")
            
            if st.session_state.current_q < len(st.session_state.quiz_data)-1:
                st.session_state.current_q += 1
                st.rerun()
            else:
                st.session_state.quiz_state = "end"
                st.rerun()

    elif st.session_state.quiz_state == "end":
        st.title(f"Score: {st.session_state.score}")
        if st.button("Restart"): 
            st.session_state.quiz_state = "setup"
            st.rerun()

# === FORMULA CARDS ===
elif mode == "🧮 Formula Cards":
    topic = st.text_input("Topic:")
    if st.button("Generate") and topic:
        st.markdown(get_gemini_response(f"Formula sheet for {topic}. Table format."))
        add_xp(20)

# === DOCUMENT READER ===
elif mode == "📄 Document Reader":
    f = st.file_uploader("Upload PDF/PPTX")
    if f:
        txt = extract_text(f)
        q = st.text_input("Question:")
        if st.button("Ask"): st.write(get_gemini_response(f"Doc: {txt}\nQ: {q}"))

# === GENERAL TOOLS (Chat, Vision, YouTube, Search) ===
elif mode == "💬 Chat":
    user_input = st.chat_input("Chat...")
    if user_input: st.write(get_gemini_response(user_input))

elif mode == "🖼️ Vision":
    f = st.file_uploader("Img")
    if f: 
        img = Image.open(f)
        st.image(img, width=200)
        if st.button("Analyze"): st.write(get_gemini_response(model.generate_content(["Describe", img]).text))

elif mode == "📺 YouTube Summarizer":
    url = st.text_input("URL")
    if url and st.button("Summarize"):
        try:
            vid = url.split("v=")[1].split("&")[0] if "v=" in url else url.split("/")[-1]
            txt = YouTubeTranscriptApi.get_transcript(vid)
            st.write(get_gemini_response(f"Summarize: {' '.join([t['text'] for t in txt])}"))
        except: st.error("No captions.")

elif mode == "🌐 Web Search":
    q = st.text_input("Search")
    if q and st.button("Go"):
        try: st.write(genai.GenerativeModel("gemini-2.5-flash", tools='google_search_retrieval').generate_content(q).text)
        except: st.write(get_gemini_response(q))